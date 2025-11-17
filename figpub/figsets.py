import os
import sys
import subprocess
import tempfile
from matplotlib import pyplot as plt
import numpy as np

# --- PPT 리포트 생성을 위한 라이브러리 ---
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
except ImportError:
    print("='python-pptx' 라이브러리가 필요합니다. 'pip install python-pptx'를 실행해주세요.")
    # 이 경우, create_report는 작동하지 않습니다.


class PanelChild:
    """Individual panel within a publication figure"""
    
    def __init__(self, parent_fig, lbwh, label=None, comment='...', draw=None):
        """
        Parameters
        ----------
        parent_fig : PubFig
            Parent figure object
        lbwh : array-like
            [left, bottom, width, height] in width-normalized units
        label : str, optional
            Label for the panel
        draw : callable, optional
            Function(ax) to draw on this panel
        """
        self.parent = parent_fig
        self.lbwh = np.asarray(lbwh, dtype=float)
        self.label = label if label is not None else 'unknown'
        self.comment = comment
        self.draw = draw
        self.ax = None         # <--- 지연 초기화

    @property
    def left(self):
        return self.lbwh[0]
    
    @property
    def bottom(self):
        return self.lbwh[1]
    
    @property
    def width(self):
        return self.lbwh[2]
    
    @property
    def height(self):
        return self.lbwh[3]

    def render(self):
        """
        (PubFig.render()에 의해 호출됨)
        실제 Matplotlib Axes 객체를 생성합니다.
        """
        if self.ax is None:
            if self.parent.fig is None:
                self.parent.render()
                
            # --- [수정됨] ---
            # 원본의 올바른 정규화 로직으로 복원합니다.
            # (lbwh는 'width' 기준 상대 좌표, add_axes는 'figure' 기준 상대 좌표)
            # b_fig = (b_w * W_px) / H_px = b_w / height_u
            lbwh_axu = self.lbwh * np.array([1, 1/self.parent.height_u, 1, 1/self.parent.height_u])
            
            self.ax = self.parent.fig.add_axes(lbwh_axu)
            # --- [수정 완료] ---

    def resize(self, width=None, height=None, anchor='bottom_left'):
        anchor_point = self.get_point(anchor)
        
        if width is not None:
            self.lbwh[2] = width
        if height is not None:
            self.lbwh[3] = height
        
        new_anchor_point = self.get_point(anchor)
        
        self.lbwh[0] += (anchor_point[0] - new_anchor_point[0])
        self.lbwh[1] += (anchor_point[1] - new_anchor_point[1])
        
        self._update_axes()

    def reduce(self, w_reduce=0, h_reduce=0, anchor='bottom_left'):
        new_width = self.width - w_reduce
        new_height = self.height - h_reduce
        self.resize(width=new_width, height=new_height, anchor=anchor)
    
    def translate(self, dx=0, dy=0):
        self.lbwh[0] += dx
        self.lbwh[1] += dy
        self._update_axes()

    def set_position(self, left=None, bottom=None):
        if left is not None:
            self.lbwh[0] = left
        if bottom is not None:
            self.lbwh[1] = bottom
        self._update_axes()
    
    def get_point(self, anchor='center'):
        if isinstance(anchor, tuple):
            x_frac, y_frac = anchor
            return np.array([
                self.left + self.width * x_frac,
                self.bottom + self.height * y_frac
            ])
        
        anchor_map = {
            'center': (0.5, 0.5), 'top': (0.5, 1.0), 'bottom': (0.5, 0.0),
            'left': (0.0, 0.5), 'right': (1.0, 0.5), 'top_left': (0.0, 1.0),
            'top_right': (1.0, 1.0), 'bottom_left': (0.0, 0.0), 'bottom_right': (1.0, 0.0),
        }
        
        if anchor not in anchor_map:
            raise ValueError(f"Unknown anchor: {anchor}")
        
        return self.get_point(anchor_map[anchor])
    
    def _update_axes(self):
        """Update axes position after transformation"""
        if self.ax is not None:
            
            # --- [수정됨] ---
            # 원본의 올바른 정규화 로직으로 복원합니다.
            lbwh_axu = self.lbwh * np.array([1, 1/self.parent.height_u, 1, 1/self.parent.height_u])
            self.ax.set_position(lbwh_axu)
            # --- [수정 완료] ---
    
    def plot_layout(self):
        """Draw layout helper (cross and label)"""
        if self.ax is None:
            self.render()
            
        self.ax.set_xticks([])
        self.ax.set_yticks([])
        self.ax.plot([0, 1], [0, 1], 'k', linewidth=0.1)
        self.ax.plot([1, 0], [0, 1], 'k', linewidth=0.1)
        self.ax.text(0, 1, self.label, transform=self.ax.transAxes, 
                     ha='left', va='top', fontsize=12)
        self.ax.text(.5, .5, self.comment, transform=self.ax.transAxes, 
                     ha='center', va='center', fontsize=6)
        self.ax.set_xlim(0, 1)
        self.ax.set_ylim(0, 1)

    def plot_draw(self):
        """
        지정된 'draw' 함수를 사용해 실제 데이터를 플롯합니다.
        draw 함수가 없으면 layout 헬퍼를 대신 그립니다.
        """
        if self.ax is None:
            self.render()
            
        if self.draw is not None:
            self.draw(self.ax)
        else:
            self.plot_layout()

 
class PubFig:
    MM_PER_INCH = 25.4
    WIDTH_2COL = 178
    WIDTH_1COL = 86
    
    def __init__(self, width, height_u, width_rescale=1,
                 figtitle="Untitled Figure", 
                 keyword_info=None, 
                 keyword_argument=None):
        
        if width == '1col':
            self.width_pure = PubFig.WIDTH_1COL
        elif width == '2col':
            self.width_pure = PubFig.WIDTH_2COL
        else:
            self.width_pure = width
            
        self.width = self.width_pure * width_rescale # <--- 최종 너비 (mm)
        self.height_u = height_u                     # <--- 높이/너비 비율
        
        # --- 리포트용 메타데이터 ---
        self.figtitle = figtitle
        self.keyword_info = keyword_info if keyword_info is not None else []
        self.keyword_argument = keyword_argument if keyword_argument is not None else []
        
        # --- 지연 초기화 ---
        self.fig = None        # <--- fig를 None으로 초기화
        self.fignum = None     # <--- fignum도 None으로 초기화
        self.children = []

    @staticmethod
    def mm_to_inch(mm):
        return mm / PubFig.MM_PER_INCH
    
    @property
    def height(self):
        # <--- 높이 (mm)
        return self.width * self.height_u
    
    def render(self):
        """
        실제 Matplotlib Figure 객체를 생성하고,
        모든 자식 패널의 렌더링을 트리거합니다.
        """
        if self.fig is None: # <--- 아직 렌더링되지 않았다면
            figsize = (PubFig.mm_to_inch(self.width), 
                       PubFig.mm_to_inch(self.height))
            self.fig = plt.figure(figsize=figsize)
            self.fignum = self.fig.number
            
            # 모든 자식 패널도 렌더링
            for child in self.children:
                child.render()
    def close(self):
        """
        Matplotlib Figure를 닫고,
        Figure와 모든 자식 Panel의 Axes 참조를 None으로 리셋합니다.
        이를 통해 객체를 재사용(재-렌더링)할 수 있게 됩니다.
        """
        if self.fig is not None:
            # 1. Matplotlib 백엔드에서 창을 닫아 메모리 해제
            plt.close(self.fig)
            
            # 2. 부모(Figure) 참조 리셋
            self.fig = None
            self.fignum = None
            
            # 3. [핵심] 모든 자식(PanelChild)의 Axes 참조 리셋
            for child in self.children:
                child.ax = None
    
    def add_child(self, lbwh=None, label=None, anchor=None, xy=None, wh=None, comment='...', draw=None):
        if lbwh is not None:
            child = PanelChild(self, lbwh, label, comment=comment, draw=draw)
        elif anchor is not None and xy is not None and wh is not None:
            xy = np.asarray(xy)
            wh = np.asarray(wh)
            
            if isinstance(anchor, tuple):
                x_frac, y_frac = anchor
            else:
                anchor_map = {
                    'center': (0.5, 0.5), 'top': (0.5, 1.0), 'bottom': (0.5, 0.0),
                    'left': (0.0, 0.5), 'right': (1.0, 0.5), 'top_left': (0.0, 1.0),
                    'top_right': (1.0, 1.0), 'bottom_left': (0.0, 0.0), 'bottom_right': (1.0, 0.0),
                }
                if anchor not in anchor_map:
                    raise ValueError(f"Unknown anchor: {anchor}")
                x_frac, y_frac = anchor_map[anchor]
            
            left = xy[0] - wh[0] * x_frac
            bottom = xy[1] - wh[1] * y_frac
            
            lbwh = [left, bottom, wh[0], wh[1]]
            child = PanelChild(self, lbwh, label, comment=comment, draw=draw)
        else:
            raise ValueError("Either provide 'lbwh' or all of 'anchor', 'xy', and 'wh'")
        
        self.children.append(child)
        return child
    
    def get_child(self, identifier):
        if isinstance(identifier, int):
            return self.children[identifier]
        elif isinstance(identifier, str):
            for child in self.children:
                if child.label == identifier:
                    return child
            raise ValueError(f"No child found with label: {identifier}")
        else:
            raise TypeError(f"identifier must be int or str, not {type(identifier)}")
    
    def plot_layout(self):
        """Plot layout for all children"""
        if self.fig is None: # <--- 렌더링이 필요하면 자동 렌더링
            self.render()
            
        for child in self.children:
            child.plot_layout()
        

class PubProject:
    """
    여러 개의 PubFig 객체를 하나의 프로젝트로 관리합니다.
    """
    
    def __init__(self, *figs, title="Untitled Project", synopsis=None):
        self.figs = []
        for fig in figs:
            if not isinstance(fig, PubFig):
                raise TypeError(f"모든 인자는 PubFig 객체여야 합니다. {type(fig)} 타입이 입력되었습니다.")
            self.figs.append(fig)
            
        # --- 리포트용 메타데이터 ---
        self.title = title
        # synopsis는 4개 문자열 리스트를 가정
        self.synopsis = synopsis if synopsis is not None else [""] * 4
        
        print(f"PubProject '{self.title}'가 {len(self.figs)}개의 Figure로 생성되었습니다.")

    def __getitem__(self, index):
        return self.figs[index]

    def __len__(self):
        return len(self.figs)

    def plot_layouts(self):
        """프로젝트에 포함된 모든 Figure의 레이아웃을 그립니다."""
        print("레이아웃 플로팅 시작...")
        for i, pub_fig in enumerate(self.figs):
            try:
                pub_fig.render() # <--- 여기서 명시적으로 렌더링
                pub_fig.plot_layout()
                
                # 창 위치 이동 (렌더링 후에만 가능)
                self._move_window(pub_fig, i)
            except Exception as e:
                print(f"Figure {i} 레이아웃 플롯팅 중 오류 발생: {e}")

    def plot_draws(self):
        """프로젝트에 포함된 모든 Figure의 실제 데이터 플롯을 그립니다."""
        print("데이터 플로팅 시작...")
        for i, pub_fig in enumerate(self.figs):
            try:
                pub_fig.render() # <--- 여기서 명시적으로 렌더링
                
                # 모든 자식의 plot_draw 호출
                for child in pub_fig.children:
                    child.plot_draw()
                    
                # 창 위치 이동 (렌더링 후에만 가능)
                self._move_window(pub_fig, i)
            except Exception as e:
                print(f"Figure {i} 플롯팅 중 오류 발생: {e}")

    def _move_window(self, pub_fig, index):
        """Helper to move the figure window."""
        try:
            # 백엔드에 따라 작동하지 않을 수 있음 (e.g., inline)
            manager = pub_fig.fig.canvas.manager
            x_pos = 500 + index * 200
            y_pos = 100 + index * 50
            manager.window.move(x_pos, y_pos)
        except Exception:
            pass # GUI 백엔드가 아니면 조용히 실패

    def show(self):
        """plt.show()를 호출하여 모든 Figure를 한꺼번에 보여줍니다."""
        print("plt.show() 호출. (모든 창 닫기 전까지 스크립트 대기)")
        plt.show()

    def save_all(self, directory='.', prefix='Fig', format='pdf', dpi=300, **kwargs):
        """
        프로젝트의 모든 Figure를 지정된 디렉토리에 저장합니다.
        """
        if not os.path.exists(directory):
            os.makedirs(directory, exist_ok=True)
            print(f"디렉토리를 생성했습니다: {directory}")
        
        for i, pub_fig in enumerate(self.figs):
            filename = os.path.join(directory, f"{prefix}{i+1}.{format}")
            try:
                pub_fig.render() # <--- 저장 직전에 렌더링
                
                # plot_draws를 명시적으로 호출해야 그림이 그려짐
                for child in pub_fig.children:
                    child.plot_draw()

                pub_fig.fig.savefig(filename, format=format, dpi=dpi, **kwargs)
                print(f"저장 완료: {filename}")
                
                # 저장 후 창을 닫아 메모리 관리
                plt.close(pub_fig.fig) 
                
            except Exception as e:
                print(f"Figure {i} 저장 중 오류 발생 ({filename}): {e}")

    def _open_file(self, filepath):
        """운영체제에 맞춰 생성된 파일을 엽니다."""
        try:
            if sys.platform == "win32":
                os.startfile(filepath)
            elif sys.platform == "darwin": # macOS
                subprocess.call(["open", filepath])
            else: # linux
                subprocess.call(["xdg-open", filepath])
            print(f"리포트 생성 완료: {filepath} (여는 중)")
        except Exception as e:
            print(f"리포트 파일을 여는 데 실패했습니다: {e}")

    # --- 🚀 새로운 리포트 생성기 ---
    def create_report(self, filename="report.pptx"):
        """
        PPT 리포트를 생성하고 엽니다.
        """
        
        # --- 0. 라이브러리 확인 ---
        try:
            Presentation
        except NameError:
            print("'python-pptx'가 설치되지 않아 리포트를 생성할 수 없습니다.")
            return

        # --- 1. PPT 객체 생성 (A4) ---
        prs = Presentation()
        prs.slide_height = Inches(11.69)
        prs.slide_width = Inches(8.27)
        
        # --- 2. 이미지 사전 렌더링 (임시 폴더) ---
        with tempfile.TemporaryDirectory() as temp_dir:
            print(f"리포트 생성 시작... 임시 디렉토리: {temp_dir}")
            thumbnail_pathss = []
            for i_ld in range(2):
                thumbnail_paths = []                        
                for i, pub_fig in enumerate(self.figs):
                    # 렌더링 및 플로팅
                    pub_fig.render()
                    for child in pub_fig.children:                        
                        if i_ld == 0:
                            child.plot_layout()
                        elif i_ld == 1:
                            child.plot_draw()
                        # child.plot_layout()  # 썸네일은 레이아웃만
                    
                    # 임시 파일로 저장
                    thumb_path = os.path.join(temp_dir, f"data_fig_{i+1}_{i_ld}.png")
                    pub_fig.fig.savefig(thumb_path, dpi=96, bbox_inches='tight')
                    print(pub_fig.fig)
                    # print(pub_fig.fig is None)
                    # plt.close(pub_fig.fig)  # <--- 메모리 해제
                    # print(pub_fig.fig is None)
                    # pub_fig.fig = None  # <--- 메모리 해제
                    pub_fig.close()  # <--- 메모리 해제
                    thumbnail_paths.append(thumb_path)
                    
                    # [중요] 메모리에서 창 닫기
                    
                thumbnail_pathss.append(thumbnail_paths)
            print(thumbnail_pathss)
            
            print(f"{len(thumbnail_paths)}개의 피겨 썸네일 생성 완료.")

            # --- 3. 슬라이드 1: 프로젝트 요약 ---
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6: Blank
            
            # --- 3a. 왼쪽: 제목 및 개요 ---
            left_tx = Inches(0.5)
            top_tx = Inches(0.5)
            width_tx = Inches(3.5)
            
            # 제목
            title_box = slide.shapes.add_textbox(left_tx, top_tx, width_tx, Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = self.title
            p.font.bold = True
            p.font.size = Pt(24)
            
            # 개요
            syn_box = slide.shapes.add_textbox(left_tx, top_tx + Inches(1.2), width_tx, Inches(4))
            tf = syn_box.text_frame
            for line in self.synopsis:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)
            
            # --- 3b. 오른쪽: A4 프레임 및 썸네일 ---
            # frame_left = Inches(4.2)
            # frame_top = Inches(1.5)
            # frame_width = Inches(3.5)
            # frame_height = frame_width * 1.414 # A4 비율
            
            # # A4 프레임
            # frame = slide.shapes.add_shape(
            #     MSO_SHAPE.RECTANGLE, frame_left, frame_top, frame_width, frame_height
            # )
            # frame.fill.background() # 채우기 없음
            # frame.line.color.rgb = RGBColor(0, 0, 0) # 검은색 테두리

            # # 썸네일 (2x2 그리드)
            # thumb_w = (frame_width / 2) - Inches(0.1)
            # thumb_h = (frame_height / 2) - Inches(0.1)
            # positions = [
            #     (frame_left + Inches(0.05), frame_top + Inches(0.05)), # (0, 0)
            #     (frame_left + thumb_w + Inches(0.15), frame_top + Inches(0.05)), # (0, 1)
            #     (frame_left + Inches(0.05), frame_top + thumb_h + Inches(0.15)), # (1, 0)
            #     (frame_left + thumb_w + Inches(0.15), frame_top + thumb_h + Inches(0.15)) # (1, 1)
            # ]
            
            # for i, thumb_path in enumerate(thumbnail_paths):
            #     if i >= 4: break # 최대 4개
            #     slide.shapes.add_picture(thumb_path, positions[i][0], positions[i][1], width=thumb_w)

            # --- 4. 슬라이드 2 ~ (N+1): 피겨 상세 ---
            MARGIN = Inches(0.5)
            SLIDE_WIDTH = prs.slide_width
            ONE_COL_IMG_WIDTH = Inches(3.5)
            TWO_COL_IMG_WIDTH = SLIDE_WIDTH - 2 * MARGIN

            for i, pub_fig in enumerate(self.figs):
                
                for ii in range(2):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    thumb_path = thumbnail_pathss[ii][i]
                    
                    # 원본 비율 계산 (높이/너비)
                    aspect_ratio = pub_fig.height_u # 원본 종횡비 사용

                    # [핵심 조건문]
                    if pub_fig.width_pure > PubFig.WIDTH_1COL:
                        # 2단 컬럼 (중앙 정렬)
                        display_width = TWO_COL_IMG_WIDTH
                        display_width = Cm(pub_fig.width / 10)  # mm -> cm                    
                        display_height = Cm(pub_fig.height / 10)  # mm -> cm
                        
                        left = (SLIDE_WIDTH - display_width) / 2 # 중앙
                        top = MARGIN*2
                        slide.shapes.add_picture(thumb_path, left, top, width=display_width)
                        
                        # 텍스트 위치 (이미지 아래)
                        text_top = top + display_width + Inches(0.2)
                        text_left = MARGIN
                        text_width = TWO_COL_IMG_WIDTH
                        text_height = Inches(4)

                    else:
                        # 1단 컬럼 (왼쪽 정렬)
                        # display_width = ONE_COL_IMG_WIDTH
                        # display_height = display_width * aspect_ratio * (pub_fig.width / pub_fig.width_pure) # 스케일링 보정
                        display_width = Cm(pub_fig.width / 10)  # mm -> cm                    
                        display_height = Cm(pub_fig.height / 10)  # mm -> cm
                        left = MARGIN # 왼쪽
                        top = MARGIN*2
                        slide.shapes.add_picture(thumb_path, left, top, width=display_width)
                        
                        # 텍스트 위치 (오른쪽)
                        text_top = top + display_width*2 + Inches(0.2)
                        text_left = MARGIN
                        text_width = SLIDE_WIDTH - text_left - MARGIN
                        text_height = Inches(10)

                    # 텍스트박스 추가
                    txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                    tf = txBox.text_frame
                    
                    # Figtitle
                    p = tf.paragraphs[0]
                    p.text = f"Figure {i+1}: {pub_fig.figtitle}"
                    p.font.bold = True
                    p.font.size = Pt(14)
                    
                    # Keyword Info
                    p = tf.add_paragraph()
                    p.text = "Info:"
                    p.font.bold = True
                    p.font.size = Pt(11)
                    for info in pub_fig.keyword_info:
                        p = tf.add_paragraph()
                        p.text = info
                        p.level = 1 # 들여쓰기
                    
                    # Keyword Argument
                    p = tf.add_paragraph()
                    p.text = "Argument:"
                    p.font.bold = True
                    p.font.size = Pt(11)
                    for arg in pub_fig.keyword_argument:
                        p = tf.add_paragraph()
                        p.text = arg
                        p.level = 1

            # --- 5. 저장 및 열기 ---
            prs.save(filename)
            self._open_file(filename)


# -----------------------------------------------------------------
# --- 사용 예시 ---
# -----------------------------------------------------------------
if __name__ == '__main__':

    # --- 1. 플로팅 함수 정의 ---
    def draw_schematic(ax):
        """(a) Schematic of DWN"""
        ax.set_title("Schematic of DWN", fontsize=10)
        ax.plot(np.linspace(0, 10, 100), np.sin(np.linspace(0, 10, 100)), 'b-')
        ax.set_xticks([])
        ax.set_yticks([])
        ax.text(0.5, 0.5, "Schematic Data", transform=ax.transAxes, ha='center')

    def draw_topo(ax):
        """(b) Topology Data"""
        ax.set_title("STM Topography", fontsize=10)
        ax.imshow(np.random.rand(20, 20), cmap='viridis')
        ax.text(0.1, 0.1, "Scalebar 10nm", transform=ax.transAxes, color='white')

    def draw_fft(ax):
        """(c) FFT"""
        ax.set_title("FFT", fontsize=10)
        data = np.zeros((10, 10))
        data[2, 2] = 1
        data[8, 8] = 1
        ax.imshow(data, cmap='hot')
    
    def draw_graph(ax):
        """(d) Line Profile"""
        ax.set_title("Line Profile", fontsize=10)
        ax.plot([1, 2, 3, 4], [10, 5, 8, 12], 'r-o')
        ax.set_xlabel("Voltage (V)")
        ax.set_ylabel("dI/dV (a.u.)")

    # --- 2. PubFig 객체 생성 (메타데이터 포함) ---
    
    # Figure 1 (2단 컬럼)
    fig1 = PubFig(
        width='2col', height_u=0.4,width_rescale=0.6, # <--- 2col, 0.4 비율
        figtitle="Overview of Domain Wall Network",
        keyword_info=["STM Data (78K)", "T=TaSe2"],
        keyword_argument=["Shows the overall domain structure.", "Confirms 3Q CDW."]
    )
    # lbwh는 (너비 1 기준) 상대 좌표
    fig1.add_child([0, 0, 0.5, 1], label='a', comment='Schematic', draw=draw_schematic)
    fig1.add_child([0.5, 0, 0.5, 1], label='b', comment='Topo', draw=draw_topo)
    
    # Figure 2 (1단 컬럼)
    fig2 = PubFig(
        width='1col', height_u=1.0, # <--- 1col, 1.0 비율 (정사각형)
        figtitle="FFT Analysis",
        keyword_info=["Fig 1b data", "High-pass filter"],
        keyword_argument=["Q-peaks clearly visible.", "No satellite peaks observed."]
    )
    fig2.add_child([0, 0, 1, 1], label='a', comment='FFT', draw=draw_fft)

    # Figure 3 (1단 컬럼)
    fig3 = PubFig(
        width='1col', height_u=0.6, # <--- 1col, 0.6 비율
        figtitle="Spectroscopy",
        keyword_info=["V_bias = -100mV", "I_set = 100pA"],
        keyword_argument=["Gap observed at Fermi level.", "Consistent with CDW phase."]
    )
    fig3.add_child([0, 0, 1, 1], label='a', comment='dIdV', draw=draw_graph)
    
    # Figure 4 (1단 컬럼, 레이아웃만)
    fig4 = PubFig(
        width='1col', height_u=0.8,
        figtitle="Placeholder Layout",
        keyword_info=["..."],
        keyword_argument=["..."]
    )
    fig4.add_child([0.1, 0.1, 0.8, 0.8], label='a', comment='No draw func')


    # --- 3. PubProject 생성 (메타데이터 포함) ---
    my_paper = PubProject(
        fig1, fig2, fig3, fig4,
        title="CDW Solitons in 2H-TaSe2",
        synopsis=[
            "We investigated the Charge Density Wave (CDW) in 2H-TaSe2.",
            "Topological domain walls (DW) and solitons were observed.",
            "A network of these domain walls forms a 'Domain Wall Network' (DWN).",
            "Spectroscopy confirms the insulating nature of the C-phase."
        ]
    )

    # --- 4. 실행 (하나만 골라서 주석 해제) ---
    
    # 옵션 A: 화면으로 레이아웃 청사진만 보기
    # my_paper.plot_layouts()
    # my_paper.show()

    # 옵션 B: 화면으로 실제 데이터 플롯 보기
    # my_paper.plot_draws()
    # my_paper.show()
    
    # 옵션 C: PDF로 저장하기 (데이터 포함)
    # my_paper.save_all(directory="paper_figures_pdf", format="pdf")
    
    # 옵션 D: PPT 리포트 생성하기
    my_paper.create_report(filename="My_Paper_Report.pptx")